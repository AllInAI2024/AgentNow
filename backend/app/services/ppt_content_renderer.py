import os
import io
import json
import logging
from datetime import datetime
from typing import Optional, List, Dict, Any
from dataclasses import dataclass, field
from enum import Enum

logger = logging.getLogger(__name__)


class SlideLayoutType(str, Enum):
    TITLE_SLIDE = "title_slide"
    TITLE_AND_CONTENT = "title_and_content"
    TWO_CONTENT = "two_content"
    COMPARISON = "comparison"
    TITLE_ONLY = "title_only"
    BLANK = "blank"


@dataclass
class ThemeColors:
    primary: str = "#165DFF"
    secondary: str = "#4080FF"
    accent: str = "#0FC6C2"
    background: str = "#FFFFFF"
    text_primary: str = "#1D2129"
    text_secondary: str = "#4E5969"
    success: str = "#00B42A"
    warning: str = "#FF7D00"
    danger: str = "#F53F3F"


@dataclass
class Theme:
    name: str
    colors: ThemeColors
    title_font: str = "微软雅黑"
    body_font: str = "微软雅黑"
    title_font_size: int = 44
    body_font_size: int = 20


class ThemePresets:
    BLUE_THEME = Theme(
        name="蓝色商务",
        colors=ThemeColors(
            primary="#165DFF",
            secondary="#4080FF",
            accent="#0FC6C2",
        )
    )
    
    ORANGE_THEME = Theme(
        name="橙色活力",
        colors=ThemeColors(
            primary="#FF7D00",
            secondary="#FF9A2E",
            accent="#F7BA1E",
        )
    )
    
    GREEN_THEME = Theme(
        name="绿色环保",
        colors=ThemeColors(
            primary="#00B42A",
            secondary="#23C343",
            accent="#0FC6C2",
        )
    )
    
    DARK_THEME = Theme(
        name="深色科技",
        colors=ThemeColors(
            primary="#165DFF",
            secondary="#4080FF",
            background="#1D2129",
            text_primary="#FFFFFF",
            text_secondary="#C9CDD4",
        )
    )
    
    PURPLE_THEME = Theme(
        name="紫色优雅",
        colors=ThemeColors(
            primary="#722ED1",
            secondary="#9454D5",
            accent="#CBECFC",
        )
    )


class PPTContentRenderer:
    """
    PPT 内容渲染器 - 使用 python-pptx 创建美观的 PPT 文件
    """
    
    THEME_MAP = {
        "蓝色商务": ThemePresets.BLUE_THEME,
        "正式汇报": ThemePresets.BLUE_THEME,
        "标准模板": ThemePresets.BLUE_THEME,
        "公司标准": ThemePresets.BLUE_THEME,
        "默认": ThemePresets.BLUE_THEME,
        "橙色活力": ThemePresets.ORANGE_THEME,
        "销售展示": ThemePresets.ORANGE_THEME,
        "客户介绍": ThemePresets.GREEN_THEME,
        "绿色环保": ThemePresets.GREEN_THEME,
        "深色科技": ThemePresets.DARK_THEME,
        "紫色优雅": ThemePresets.PURPLE_THEME,
    }
    
    def __init__(self, theme_name: Optional[str] = None):
        self.theme = self._get_theme(theme_name)
        self._ensure_pptx_installed()
    
    def _get_theme(self, theme_name: Optional[str]) -> Theme:
        if not theme_name:
            return ThemePresets.BLUE_THEME
        
        for key, theme in self.THEME_MAP.items():
            if key in theme_name or theme_name in key:
                return theme
        
        return ThemePresets.BLUE_THEME
    
    def _ensure_pptx_installed(self) -> bool:
        try:
            import pptx
            return True
        except ImportError:
            logger.warning("python-pptx is not installed. PPT generation may not work.")
            return False
    
    def _create_title_slide(
        self,
        prs,
        title: str,
        subtitle: Optional[str] = None,
    ):
        """
        创建封面页
        """
        title_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(title_slide_layout)
        
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        left = Inches(0)
        top = Inches(0)
        width = Inches(10)
        height = Inches(5.625)
        
        bg_shape = slide.shapes.add_shape(1, left, top, width, height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor.from_string(self.theme.colors.primary.lstrip('#'))
        
        title_left = Inches(0.5)
        title_top = Inches(2)
        title_width = Inches(9)
        title_height = Inches(1.2)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(self.theme.title_font_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = self.theme.title_font
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_top = Inches(3.3)
            subtitle_height = Inches(0.8)
            
            subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, subtitle_height)
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.size = Pt(24)
            p_sub.font.color.rgb = RGBColor(230, 230, 230)
            p_sub.font.name = self.theme.body_font
            p_sub.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_title_and_content_slide(
        self,
        prs,
        title: str,
        subtitle: Optional[str] = None,
        bullets: Optional[List[str]] = None,
        speaker_notes: Optional[str] = None,
    ):
        """
        创建标题和内容页
        """
        title_content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(title_content_layout)
        
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            title_para = slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor.from_string(self.theme.colors.primary.lstrip('#'))
            title_para.font.name = self.theme.title_font
        
        body_shape = None
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                body_shape = shape
                break
        
        if body_shape:
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            if subtitle:
                tf.text = subtitle
                first_para = tf.paragraphs[0]
                first_para.font.size = Pt(self.theme.body_font_size)
                first_para.font.color.rgb = RGBColor.from_string(self.theme.colors.text_primary.lstrip('#'))
                first_para.font.name = self.theme.body_font
            
            bullets = bullets or []
            for bullet in bullets:
                if bullet:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.font.size = Pt(self.theme.body_font_size)
                    p.font.color.rgb = RGBColor.from_string(self.theme.colors.text_primary.lstrip('#'))
                    p.font.name = self.theme.body_font
                    p.level = 0
            
            if speaker_notes:
                p = tf.add_paragraph()
                p.text = ""
                p = tf.add_paragraph()
                p.text = f"【演讲备注】{speaker_notes}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor.from_string(self.theme.colors.text_secondary.lstrip('#'))
                p.font.italic = True
        
        return slide
    
    def _create_section_header_slide(
        self,
        prs,
        title: str,
        subtitle: Optional[str] = None,
    ):
        """
        创建章节标题页
        """
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        left = Inches(0)
        top = Inches(0)
        width = Inches(10)
        height = Inches(5.625)
        
        bg_shape = slide.shapes.add_shape(1, left, top, width, height)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
        
        title_left = Inches(0.5)
        title_top = Inches(2)
        title_width = Inches(9)
        title_height = Inches(1.0)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(self.theme.colors.primary.lstrip('#'))
        p.font.name = self.theme.title_font
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_top = Inches(3.2)
            subtitle_height = Inches(0.6)
            
            subtitle_box = slide.shapes.add_textbox(title_left, subtitle_top, title_width, subtitle_height)
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = RGBColor.from_string(self.theme.colors.text_secondary.lstrip('#'))
            p_sub.font.name = self.theme.body_font
            p_sub.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_comparison_slide(
        self,
        prs,
        title: str,
        bullets: Optional[List[str]] = None,
    ):
        """
        创建对比页
        """
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(9)
        title_height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor.from_string(self.theme.colors.primary.lstrip('#'))
        
        bullets = bullets or []
        num_bullets = len(bullets)
        
        if num_bullets > 0:
            half = num_bullets // 2
            left_bullets = bullets[:half]
            right_bullets = bullets[half:]
            
            left_col_left = Inches(0.5)
            right_col_left = Inches(5.25)
            col_top = Inches(1.2)
            col_width = Inches(4.25)
            col_height = Inches(4)
            
            self._add_bullet_column(
                slide, left_col_left, col_top, col_width, col_height,
                left_bullets, "左侧"
            )
            
            self._add_bullet_column(
                slide, right_col_left, col_top, col_width, col_height,
                right_bullets, "右侧"
            )
        
        return slide
    
    def _add_bullet_column(
        self,
        slide,
        left,
        top,
        width,
        height,
        bullets: List[str],
        label: str,
    ):
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(bullet)
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor.from_string(self.theme.colors.text_primary.lstrip('#'))
            p.font.name = self.theme.body_font
            p.level = 0
    
    def render_slide(
        self,
        prs,
        slide_data: Dict[str, Any],
    ):
        """
        根据 slide 数据渲染单页 PPT
        """
        layout = slide_data.get("layout", "title_and_content")
        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle")
        bullets = slide_data.get("bullets", [])
        speaker_notes = slide_data.get("speaker_notes")
        
        if layout == SlideLayoutType.TITLE_SLIDE:
            return self._create_title_slide(prs, title, subtitle)
        elif layout == SlideLayoutType.TWO_CONTENT:
            return self._create_comparison_slide(prs, title, bullets)
        elif layout == SlideLayoutType.COMPARISON:
            return self._create_comparison_slide(prs, title, bullets)
        elif layout == SlideLayoutType.TITLE_ONLY:
            return self._create_section_header_slide(prs, title, subtitle)
        else:
            return self._create_title_and_content_slide(
                prs, title, subtitle, bullets, speaker_notes
            )
    
    def create_ppt(
        self,
        slides_data: List[Dict[str, Any]],
        title: str,
        subtitle: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> bytes:
        """
        创建完整的 PPT 文件
        
        Args:
            slides_data: 每页 slide 的数据列表
            title: PPT 主标题
            subtitle: PPT 副标题
            metadata: 元数据信息
        
        Returns:
            PPT 文件的字节内容
        """
        from pptx import Presentation
        from pptx.util import Pt, Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        prs = Presentation()
        
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)
        
        if slides_data:
            first_slide = slides_data[0]
            first_layout = first_slide.get("layout", "")
            
            if first_layout != SlideLayoutType.TITLE_SLIDE:
                full_subtitle = subtitle or ""
                if metadata:
                    if metadata.get("scene"):
                        full_subtitle += f"\n场景：{metadata['scene']}"
                    if metadata.get("audience"):
                        full_subtitle += f"\n受众：{metadata['audience']}"
                    if metadata.get("style"):
                        full_subtitle += f"\n风格：{self.theme.name}"
                
                self._create_title_slide(prs, title, full_subtitle)
            
            for slide_data in slides_data:
                if slides_data.index(slide_data) == 0 and first_layout != SlideLayoutType.TITLE_SLIDE:
                    continue
                self.render_slide(prs, slide_data)
        else:
            full_subtitle = subtitle or ""
            if metadata:
                if metadata.get("scene"):
                    full_subtitle += f"\n场景：{metadata['scene']}"
                if metadata.get("audience"):
                    full_subtitle += f"\n受众：{metadata['audience']}"
                if metadata.get("style"):
                    full_subtitle += f"\n风格：{self.theme.name}"
            
            self._create_title_slide(prs, title, full_subtitle)
            
            self._create_title_and_content_slide(
                prs,
                title="目录",
                subtitle="PPT 结构说明",
                bullets=[
                    "封面页",
                    "目录页",
                    "正文内容",
                    "总结页",
                ],
            )
        
        summary_slide_layout = prs.slide_layouts[1]
        summary_slide = prs.slides.add_slide(summary_slide_layout)
        
        if summary_slide.shapes.title:
            summary_slide.shapes.title.text = "感谢观看"
            title_para = summary_slide.shapes.title.text_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor.from_string(self.theme.colors.primary.lstrip('#'))
        
        if metadata:
            for shape in summary_slide.shapes:
                if hasattr(shape, 'text_frame') and shape != summary_slide.shapes.title:
                    tf = shape.text_frame
                    tf.text = ""
                    
                    p = tf.paragraphs[0]
                    p.text = "PPT 生成信息"
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor.from_string(self.theme.colors.text_primary.lstrip('#'))
                    
                    p = tf.add_paragraph()
                    p.text = f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor.from_string(self.theme.colors.text_secondary.lstrip('#'))
                    
                    if metadata.get("version"):
                        p = tf.add_paragraph()
                        p.text = f"版本：v{metadata['version']}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor.from_string(self.theme.colors.text_secondary.lstrip('#'))
                    
                    break
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()


def create_ppt_from_data(
    slides_data: List[Dict[str, Any]],
    title: str,
    subtitle: Optional[str] = None,
    style: Optional[str] = None,
    metadata: Optional[Dict[str, Any]] = None,
) -> bytes:
    """
    根据数据创建 PPT 的便捷函数
    """
    renderer = PPTContentRenderer(style)
    
    full_metadata = metadata or {}
    if style:
        full_metadata["style"] = style
    
    return renderer.create_ppt(
        slides_data=slides_data,
        title=title,
        subtitle=subtitle,
        metadata=full_metadata,
    )
