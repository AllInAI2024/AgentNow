import math
import json
import os
import logging
import re
from typing import Optional, Tuple, List, Dict, Any
from datetime import datetime
from sqlalchemy.orm import Session

from app.config import settings
from app.models import (
    User,
    UserAgent,
    AgentTemplate,
    AgentConversation,
    AgentGeneratedFile,
    AgentOperationLog,
)
from app.schemas import (
    UserAgentListResponse,
    UserAgentWithTemplateResponse,
    EnableAgentResponse,
    AgentConversationListResponse,
    ChatResponse,
    ConversationDetailResponse,
)
from app.services.hermes_profile_service import HermesProfileService
from app.services.knowledge_retrieval_service import KnowledgeRetrievalService
from app.services.hermes_chat_service import HermesChatService
from app.services.ppt_content_renderer import create_ppt_from_data

logger = logging.getLogger(__name__)


class AgentService:
    def __init__(self, db: Session):
        self.db = db
        self._profile_service: Optional[HermesProfileService] = None
        self._hermes_chat_service: Optional[HermesChatService] = None
        self._use_hermes: bool = self._check_hermes_available()
    
    def _check_hermes_available(self) -> bool:
        hermes_url = settings.HERMES_BASE_URL or "http://127.0.0.1:8642"
        logger.info(f"Hermes URL configured as: {hermes_url}")
        return True
    
    def _get_profile_service(self) -> HermesProfileService:
        if self._profile_service is None:
            self._profile_service = HermesProfileService()
        return self._profile_service
    
    def _get_hermes_chat_service(self) -> HermesChatService:
        if self._hermes_chat_service is None:
            self._hermes_chat_service = HermesChatService()
        return self._hermes_chat_service
    
    def _build_messages_history(
        self,
        conversation: AgentConversation,
        user_message: Optional[str]
    ) -> List[Dict[str, Any]]:
        messages = []
        history = conversation.messages_json or []
        
        for msg in history:
            if isinstance(msg, dict) and "role" in msg and "content" in msg:
                messages.append({
                    "role": msg["role"],
                    "content": msg["content"]
                })
        
        if user_message:
            messages.append({
                "role": "user",
                "content": user_message
            })
        
        return messages

    def _build_hermes_system_prompt(self, user_agent: UserAgent) -> str:
        snapshot = user_agent.config_snapshot or {}
        role_prompt = str(snapshot.get("role_prompt") or "").strip()
        system_prompt = str(snapshot.get("system_prompt") or "").strip()

        parts = [p for p in [role_prompt, system_prompt] if p]
        template_prompt = "\n\n".join(parts).strip()

        base_prompt = HermesChatService.PPT_ASSISTANT_PROMPT.strip()
        if template_prompt:
            return f"{template_prompt}\n\n{base_prompt}".strip()
        return base_prompt

    def _extract_ppt_requirements_from_user_text(
        self,
        text: str,
        existing: Optional[Dict[str, Any]],
    ) -> Dict[str, Any]:
        data: Dict[str, Any] = dict(existing or {})
        t = (text or "").strip()
        if not t:
            return data

        numbered = re.findall(r"(?m)^\s*([1-5])[\.\)、)\s]+([^\n]+)\s*$", t)
        if numbered:
            for idx_str, val in numbered:
                idx = int(idx_str)
                v = val.strip()
                if not v:
                    continue
                if idx == 1:
                    data["topic"] = v
                elif idx == 2:
                    data["audience"] = v
                elif idx == 3:
                    data["scene"] = v
                elif idx == 4:
                    m_pages = re.search(r"(\d{1,2})\s*页", v)
                    if m_pages:
                        try:
                            data["page_count"] = int(m_pages.group(1))
                        except Exception:
                            pass
                    else:
                        m_num = re.search(r"\b(\d{1,2})\b", v)
                        if m_num:
                            try:
                                data["page_count"] = int(m_num.group(1))
                            except Exception:
                                pass
                elif idx == 5:
                    data["style"] = v

        if not numbered:
            lines = [ln.strip() for ln in t.splitlines() if ln.strip()]
            has_labels = any(k in t for k in ["主题", "受众", "场景", "页数", "风格"])
            if not has_labels and len(lines) >= 5:
                data.setdefault("topic", lines[0])
                data.setdefault("audience", lines[1])
                data.setdefault("scene", lines[2])
                m_pages = re.search(r"(\d{1,2})\s*页", lines[3])
                if m_pages:
                    try:
                        data.setdefault("page_count", int(m_pages.group(1)))
                    except Exception:
                        pass
                else:
                    m_num = re.search(r"\b(\d{1,2})\b", lines[3])
                    if m_num:
                        try:
                            data.setdefault("page_count", int(m_num.group(1)))
                        except Exception:
                            pass
                data.setdefault("style", lines[4])

        def _pick(patterns: List[str]) -> Optional[str]:
            for p in patterns:
                m = re.search(p, t)
                if m:
                    val = (m.group(1) or "").strip()
                    if val:
                        return val
            return None

        topic = _pick([r"主题[:：]\s*([^\n；;]+)", r"主题是\s*([^\n；;]+)"])
        if topic:
            data["topic"] = topic

        audience = _pick([r"受众[:：]\s*([^\n；;]+)", r"给谁看[:：]?\s*([^\n；;]+)"])
        if audience:
            data["audience"] = audience

        scene = _pick([r"场景[:：]\s*([^\n；;]+)", r"用于[:：]?\s*([^\n；;]+)"])
        if scene:
            data["scene"] = scene

        style = _pick([r"风格[:：]\s*([^\n；;]+)"])
        if style:
            data["style"] = style

        m_pages = re.search(r"(?:页数[:：]\s*)?(\d{1,2})\s*页", t)
        if m_pages:
            try:
                data["page_count"] = int(m_pages.group(1))
            except Exception:
                pass

        if not str(data.get("topic") or "").strip():
            topic_candidates = [
                "公司介绍",
                "企业介绍",
                "产品介绍",
                "客户拜访汇报",
                "售前宣讲",
                "解决方案",
                "商业计划书",
                "季度汇报",
                "年终总结",
                "路演融资",
            ]
            for cand in topic_candidates:
                if cand in t:
                    data["topic"] = cand
                    break

        if not str(data.get("topic") or "").strip():
            m_topic1 = re.search(r"(?:帮我|请|麻烦)?(?:做|写|生成)(?:一份|一个)?(.+?)(?:的)?\s*(?:PPT|ppt)", t)
            if m_topic1:
                topic = (m_topic1.group(1) or "").strip("：:，,。. ")
                topic = re.sub(r"(\d{1,2})\s*页.*$", "", topic).strip("：:，,。. ")
                if topic and len(topic) <= 80 and topic not in ["PPT", "ppt", "演示", "演示文稿"]:
                    data["topic"] = topic

        if not str(data.get("topic") or "").strip():
            m_topic2 = re.search(r"(\d{1,2})\s*页\s*([^\n]{2,80}?)(?:PPT|ppt)", t)
            if m_topic2:
                topic = (m_topic2.group(2) or "").strip("：:，,。. ")
                if topic and topic not in ["PPT", "ppt"]:
                    data["topic"] = topic

        template_path = self._extract_template_path_from_message(t)
        if template_path:
            data["template_path"] = template_path
        if not str(data.get("style") or "").strip():
            normalized_style = self._normalize_template_style_from_message(t)
            if normalized_style:
                data["style"] = normalized_style

        if not str(data.get("topic") or "").strip():
            lines = [ln.strip() for ln in t.splitlines() if ln.strip()]
            if len(lines) == 1:
                candidate = lines[0]
                if candidate and len(candidate) <= 80:
                    if not re.search(r"\b(\d{1,2})\s*页\b", candidate) and "页" not in candidate:
                        data["topic"] = candidate[:80]

        return data

    def _is_ppt_requirements_complete(self, requirements: Dict[str, Any]) -> bool:
        required_keys = ["topic", "page_count"]
        for k in required_keys:
            v = requirements.get(k)
            if v is None:
                return False
            if isinstance(v, str) and not v.strip():
                return False
        return True

    def _looks_like_reasking_requirements(self, assistant_content: str) -> bool:
        c = assistant_content or ""
        keys = ["主题", "受众", "场景", "页数", "风格"]
        return sum(1 for k in keys if k in c) >= 4

    def _count_recent_reasking(self, messages_json: List[Dict[str, Any]]) -> int:
        if not messages_json:
            return 0
        count = 0
        for msg in reversed(messages_json[-10:]):
            if not isinstance(msg, dict):
                continue
            if msg.get("role") != "assistant":
                continue
            content = str(msg.get("content") or "")
            if self._looks_like_reasking_requirements(content):
                count += 1
        return count

    def _build_requirements_summary(self, requirements: Dict[str, Any]) -> str:
        keys = [
            ("topic", "主题"),
            ("audience", "受众"),
            ("scene", "场景"),
            ("page_count", "页数"),
            ("style", "风格"),
        ]
        lines: List[str] = []
        for k, label in keys:
            v = requirements.get(k)
            if v is None or v == "":
                continue
            lines.append(f"- {label}：{v}")
        if not lines:
            return ""
        return "【已知需求】\n" + "\n".join(lines)

    def _build_hermes_query(self, conversation: AgentConversation, user_message: str) -> str:
        req = self._get_conversation_requirements_data(conversation)
        summary = self._build_requirements_summary(req)
        if summary:
            return f"{summary}\n\n【用户补充】\n{user_message}"
        return user_message

    def _detect_user_confirmation(self, text: str) -> Optional[bool]:
        t = (text or "").strip().lower()
        if not t:
            return None
        confirm_keywords = ["确认", "可以", "没问题", "好的", "ok", "okay", "行", "就这样", "开始生成", "生成吧"]
        reject_keywords = ["不", "不对", "不行", "不满意", "改", "修改", "调整", "重来", "重新", "不要这样"]
        if any(k in t for k in confirm_keywords):
            return True
        if any(k in t for k in reject_keywords):
            return False
        return None

    def _extract_template_path_from_message(self, text: str) -> Optional[str]:
        t = (text or "").strip()
        if not t:
            return None
        candidates = re.findall(
            r"(?i)(file:///\\S+?\\.pptx|/\\S+?\\.pptx|[A-Za-z]:\\\\\\S+?\\.pptx|\\S+?\\.pptx)",
            t,
        )
        if not candidates:
            return None

        raw = str(candidates[0]).strip()
        raw = raw.strip('"').strip("'").strip("`").strip("<>").strip()
        raw = raw.rstrip("，,。.;；")
        if raw.lower().startswith("file:///"):
            raw = raw[len("file://") :]
        return raw or None

    def _normalize_template_style_from_message(self, text: str) -> Optional[str]:
        t = (text or "").strip()
        if not t:
            return None
        if "公司" in t and ("模板" in t or "标准" in t):
            return "公司标准"
        if "正式" in t or "商务" in t or "汇报" in t:
            return "正式汇报"
        if "客户" in t:
            return "客户介绍"
        if "销售" in t:
            return "销售展示"
        if "科技" in t:
            return "深色科技"
        if "简洁" in t or "极简" in t:
            return "默认"
        return None

    def _get_missing_requirements_for_stage(self, stage: str, requirements: Dict[str, Any]) -> List[str]:
        missing: List[str] = []
        st = (stage or "").strip() or "welcome"
        if st in ["welcome", "clarifying", "collecting_requirements"]:
            if not str(requirements.get("topic") or "").strip():
                missing.append("PPT 主题/内容")
            if requirements.get("page_count") is None:
                missing.append("页数/篇幅")
            return missing

        if st in ["template_select", "template_confirm"]:
            template_path = str(requirements.get("template_path") or "").strip()
            style = str(requirements.get("style") or "").strip()
            if not template_path and not style:
                missing.append("模板（上传/路径/或选择一种风格）")
            return missing

        return missing

    def _append_message(self, conversation: AgentConversation, role: str, content: str) -> None:
        existing = conversation.messages_json
        if isinstance(existing, list):
            items: List[Dict[str, Any]] = list(existing)
        else:
            items = []
        items.append({"role": role, "content": content})
        conversation.messages_json = list(items)
        conversation.message_count = len(items)

    def _format_outline_markdown(self, structured_result: Dict[str, Any]) -> str:
        slides = structured_result.get("slides") if isinstance(structured_result, dict) else None
        if not isinstance(slides, list) or not slides:
            return ""
        lines: List[str] = ["我已根据你的需求生成了 PPT 大纲：", ""]
        for s in slides:
            if not isinstance(s, dict):
                continue
            idx = s.get("index")
            title = s.get("title")
            bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
            if idx is None:
                lines.append(f"- {title}")
            else:
                lines.append(f"{idx}. {title}")
            for b in bullets[:6]:
                if b:
                    lines.append(f"   - {b}")
        return "\n".join(lines).strip()

    def _generate_ppt_structured_result(
        self,
        *,
        user_agent: UserAgent,
        conversation: AgentConversation,
        requirements: Dict[str, Any],
    ) -> Optional[Dict[str, Any]]:
        hermes_service = self._get_hermes_chat_service()
        system_prompt = self._build_hermes_system_prompt(user_agent)
        summary = self._build_requirements_summary(requirements)
        revision_note = str(requirements.get("revision_note") or "").strip()
        revision_text = f"\n\n【用户修改要求】\n{revision_note}" if revision_note else ""
        attempts = [
            (
                f"{summary}{revision_text}\n\n"
                "请直接输出最终结构化 JSON（type=ppt_content）。"
                "不要再提问，不要输出多余解释。"
                "输出必须只有一个 JSON 对象，不能包裹 ```。"
                "slides 里每页给出 title 与 3-5 条 bullets。"
            ).strip(),
            (
                f"{summary}{revision_text}\n\n"
                "只输出 JSON，不要任何中文解释、不要 markdown、不要代码块。"
                "JSON 结构必须满足："
                "{\"type\":\"ppt_content\",\"title\":\"...\",\"subtitle\":null,\"scene\":\"...\",\"audience\":\"...\",\"style\":\"...\",\"slides\":[{\"index\":1,\"title\":\"...\",\"bullets\":[\"...\"],\"layout\":\"title_and_content\"}]}"
            ).strip(),
        ]

        for q in attempts:
            assistant_content, stage, ppt_content, hermes_session_id = hermes_service.chat_with_ppt_assistant(
                user_message=q,
                system_prompt=system_prompt,
                profile_name=conversation.hermes_profile or user_agent.hermes_profile,
                hermes_session_id=conversation.hermes_conversation_id,
            )
            if hermes_session_id and hermes_session_id != conversation.hermes_conversation_id:
                conversation.hermes_conversation_id = hermes_session_id
            if stage:
                conversation.current_stage = stage.value

            if ppt_content:
                return {
                    "type": "ppt_content",
                    "title": ppt_content.title,
                    "subtitle": ppt_content.subtitle,
                    "scene": ppt_content.scene,
                    "audience": ppt_content.audience,
                    "style": ppt_content.style,
                    "slides": ppt_content.slides,
                }

            logger.info(
                "Hermes did not return ppt_content JSON, stage=%s, content_head=%s",
                stage.value if stage else None,
                assistant_content[:200],
            )

        return None
    
    def _chat_with_hermes(
        self,
        user_agent: UserAgent,
        conversation: AgentConversation,
        message: Optional[str],
        is_new_conversation: bool,
    ) -> Tuple[str, str, Optional[Dict[str, Any]]]:
        """
        调用 Hermes 进行对话
        
        Returns:
            (assistant_content, new_stage, structured_result)
            
        Raises:
            Exception: 如果 Hermes 服务不可用或调用失败
        """
        hermes_service = self._get_hermes_chat_service()

        if not message:
            return "", "clarifying", None

        profile_service = self._get_profile_service()
        ok, bootstrap_message = profile_service.ensure_profile_bootstrapped(
            conversation.hermes_profile or user_agent.hermes_profile
        )
        if not ok:
            raise RuntimeError(
                "Hermes Profile 未完成初始化（模型/密钥未配置）。"
                "请先在 ~/.hermes/config.yaml 选择模型，"
                "并在 ~/.hermes/.env 设置推理服务密钥（如 OPENAI_API_KEY/OPENROUTER_API_KEY），"
                f"然后重试。详情：{bootstrap_message}"
            )

        system_prompt = self._build_hermes_system_prompt(user_agent)
        query = self._build_hermes_query(conversation, message)
        logger.info(
            "Calling Hermes CLI: conversation_id=%s, is_new=%s, profile=%s, hermes_session_id=%s",
            conversation.id,
            is_new_conversation,
            conversation.hermes_profile or user_agent.hermes_profile,
            conversation.hermes_conversation_id,
        )

        assistant_content, stage, ppt_content, hermes_session_id = hermes_service.chat_with_ppt_assistant(
            user_message=query,
            system_prompt=system_prompt,
            profile_name=conversation.hermes_profile or user_agent.hermes_profile,
            hermes_session_id=conversation.hermes_conversation_id,
        )

        if hermes_session_id and hermes_session_id != conversation.hermes_conversation_id:
            conversation.hermes_conversation_id = hermes_session_id
        
        structured_result = None
        if ppt_content:
            structured_result = {
                "type": "ppt_content",
                "title": ppt_content.title,
                "subtitle": ppt_content.subtitle,
                "scene": ppt_content.scene,
                "audience": ppt_content.audience,
                "style": ppt_content.style,
                "slides": ppt_content.slides,
            }
        
        logger.info(
            "Hermes response received, stage=%s, has_content=%s, hermes_session_id=%s",
            stage,
            structured_result is not None,
            conversation.hermes_conversation_id,
        )
        
        return assistant_content, stage.value if stage else "clarifying", structured_result
    
    def _update_conversation_from_hermes(
        self,
        conversation: AgentConversation,
        user_message: Optional[str],
        assistant_content: str,
        new_stage: str,
        structured_result: Optional[Dict[str, Any]],
    ) -> None:
        """
        使用 Hermes 的结果更新会话
        """
        messages = []
        old_messages = conversation.messages_json or []
        for msg in old_messages:
            if isinstance(msg, dict) and "role" in msg and "content" in msg:
                messages.append({
                    "role": msg["role"],
                    "content": msg["content"]
                })
        
        old_count = len(messages)
        
        if user_message:
            messages.append({
                "role": "user",
                "content": user_message
            })
        
        messages.append({
            "role": "assistant",
            "content": assistant_content
        })
        
        conversation.messages_json = messages
        conversation.message_count = len(messages)
        conversation.current_stage = new_stage
        
        logger.info(
            f"Updated conversation {conversation.id}: "
            f"message_count: {old_count} -> {len(messages)}, "
            f"stage: {new_stage}, "
            f"has_structured_result: {structured_result is not None}"
        )
        logger.debug(f"Messages JSON: {json.dumps(messages, ensure_ascii=False)}")
        
        if structured_result:
            conversation.structured_result_json = structured_result
            
            if new_stage == "content_ready":
                conversation.outline_confirmed = True
                conversation.template_confirmed = True
                conversation.final_generation_confirmed = True
                conversation.status = 2
                conversation.completed_at = datetime.utcnow()
                
                requirements = conversation.requirements_json or {}
                requirements["topic"] = structured_result.get("title", requirements.get("topic", ""))
                requirements["audience"] = structured_result.get("audience", requirements.get("audience"))
                requirements["scene"] = structured_result.get("scene", requirements.get("scene"))
                requirements["style"] = structured_result.get("style", requirements.get("style"))
                conversation.requirements_json = requirements
        else:
            if new_stage == "outline_draft":
                conversation.outline_confirmed = False
            elif new_stage == "template_select":
                conversation.outline_confirmed = True
                conversation.template_confirmed = False
            elif new_stage == "final_generating":
                conversation.outline_confirmed = True
                conversation.template_confirmed = True
                conversation.final_generation_confirmed = False
    
    def get_user_agent_by_id(self, agent_id: int, user_id: int) -> Optional[UserAgent]:
        return self.db.query(UserAgent).filter(
            UserAgent.id == agent_id,
            UserAgent.user_id == user_id
        ).first()
    
    def get_user_agents(self, user_id: int) -> UserAgentListResponse:
        user_agents = self.db.query(UserAgent).filter(
            UserAgent.user_id == user_id
        ).order_by(UserAgent.last_used_at.desc()).all()
        
        items = []
        for ua in user_agents:
            template = self.db.query(AgentTemplate).filter(
                AgentTemplate.id == ua.template_id
            ).first()
            
            template_dict = None
            if template:
                template_dict = {
                    "id": template.id,
                    "code": template.code,
                    "name": template.name,
                    "description": template.description,
                    "status": template.status,
                    "is_default": template.is_default,
                    "version": template.version,
                }
            
            items.append(UserAgentWithTemplateResponse(
                id=ua.id,
                user_id=ua.user_id,
                template_id=ua.template_id,
                display_name=ua.display_name,
                hermes_profile=ua.hermes_profile,
                template_version=ua.template_version,
                config_snapshot=ua.config_snapshot,
                agent_status=ua.agent_status,
                activation_mode=ua.activation_mode,
                enabled_at=ua.enabled_at,
                last_used_at=ua.last_used_at,
                disabled_at=ua.disabled_at,
                created_at=ua.created_at,
                updated_at=ua.updated_at,
                template=template_dict,
            ))
        
        return UserAgentListResponse(
            items=items,
            total=len(items),
        )
    
    def get_user_agent_detail(
        self,
        agent_id: int,
        user_id: int,
    ) -> Tuple[Optional[UserAgentWithTemplateResponse], str]:
        user_agent = self.get_user_agent_by_id(agent_id, user_id)
        if not user_agent:
            return None, "智能体不存在或无权限访问"
        
        template = self.db.query(AgentTemplate).filter(
            AgentTemplate.id == user_agent.template_id
        ).first()
        
        template_dict = None
        if template:
            template_dict = {
                "id": template.id,
                "code": template.code,
                "name": template.name,
                "description": template.description,
                "status": template.status,
                "is_default": template.is_default,
                "version": template.version,
            }
        
        result = UserAgentWithTemplateResponse(
            id=user_agent.id,
            user_id=user_agent.user_id,
            template_id=user_agent.template_id,
            display_name=user_agent.display_name,
            hermes_profile=user_agent.hermes_profile,
            template_version=user_agent.template_version,
            config_snapshot=user_agent.config_snapshot,
            agent_status=user_agent.agent_status,
            activation_mode=user_agent.activation_mode,
            enabled_at=user_agent.enabled_at,
            last_used_at=user_agent.last_used_at,
            disabled_at=user_agent.disabled_at,
            created_at=user_agent.created_at,
            updated_at=user_agent.updated_at,
            template=template_dict,
        )
        
        return result, "获取成功"
    
    def _log_operation(
        self,
        operator_user_id: int,
        target_type: str,
        target_id: Optional[int],
        action: str,
        action_name: str,
        success: bool,
        details: Optional[Dict[str, Any]] = None,
        error_message: Optional[str] = None,
    ) -> None:
        log = AgentOperationLog(
            operator_user_id=operator_user_id,
            target_type=target_type,
            target_id=target_id,
            action=action,
            action_name=action_name,
            result_status=1 if success else 0,
            details=details,
            error_message=error_message,
            created_at=datetime.utcnow(),
        )
        self.db.add(log)
        self.db.commit()
    
    def _get_default_template(self) -> Optional[AgentTemplate]:
        return self.db.query(AgentTemplate).filter(
            AgentTemplate.is_default == True,
            AgentTemplate.status == 1,
            AgentTemplate.deleted_at.is_(None)
        ).first()
    
    def _get_template_by_id(self, template_id: int) -> Optional[AgentTemplate]:
        return self.db.query(AgentTemplate).filter(
            AgentTemplate.id == template_id,
            AgentTemplate.status == 1,
            AgentTemplate.deleted_at.is_(None)
        ).first()
    
    def _get_user_by_id(self, user_id: int) -> Optional[User]:
        return self.db.query(User).filter(
            User.id == user_id,
            User.is_active == True
        ).first()
    
    def _ensure_hermes_profile_for_user(
        self,
        user_id: int,
        tenant_id: int = 1,
    ) -> Tuple[str, bool, str]:
        profile_service = self._get_profile_service()
        
        profile_name, created, message = profile_service.ensure_profile(
            user_id=user_id,
            tenant_id=tenant_id,
        )
        
        return profile_name, created, message
    
    def enable_agent_for_user(
        self,
        user_id: int,
        template_id: Optional[int] = None,
        tenant_id: int = 1,
    ) -> Tuple[Optional[EnableAgentResponse], str]:
        user = self._get_user_by_id(user_id)
        if not user:
            return None, "用户不存在或已被停用"
        
        if template_id:
            template = self._get_template_by_id(template_id)
            if not template:
                return None, f"模板 ID={template_id} 不存在或已停用"
        else:
            template = self._get_default_template()
            if not template:
                return None, "没有可用的模板，请先配置默认模板"
        
        existing_agent = self.db.query(UserAgent).filter(
            UserAgent.user_id == user_id,
            UserAgent.template_id == template.id
        ).first()
        
        if existing_agent:
            return None, "该智能体已开通"
        
        profile_name, created_profile, profile_message = self._ensure_hermes_profile_for_user(
            user_id=user_id,
            tenant_id=tenant_id,
        )
        
        if not profile_name:
            self._log_operation(
                operator_user_id=user_id,
                target_type="user_agent",
                target_id=None,
                action="agent:enable",
                action_name="开通智能体",
                success=False,
                details={
                    "template_id": template.id,
                    "template_code": template.code,
                },
                error_message=profile_message,
            )
            return None, f"创建 Hermes Profile 失败: {profile_message}"
        
        if user.hermes_profile != profile_name:
            user.hermes_profile = profile_name
            user.updated_at = datetime.utcnow()
        
        config_snapshot = template.to_snapshot()
        
        user_agent = UserAgent(
            user_id=user_id,
            template_id=template.id,
            display_name=template.name,
            hermes_profile=profile_name,
            template_version=template.version,
            config_snapshot=config_snapshot,
            agent_status=1,
            activation_mode="auto",
            enabled_at=datetime.utcnow(),
        )
        
        self.db.add(user_agent)
        self.db.commit()
        self.db.refresh(user_agent)
        
        self._log_operation(
            operator_user_id=user_id,
            target_type="user_agent",
            target_id=user_agent.id,
            action="agent:enable",
            action_name="开通智能体",
            success=True,
            details={
                "template_id": template.id,
                "template_code": template.code,
                "hermes_profile": profile_name,
                "created_profile": created_profile,
            },
        )
        
        return EnableAgentResponse(
            user_agent=user_agent,
            created_profile=created_profile,
        ), "开通成功"
    
    def get_or_create_agent_for_user(
        self,
        user_id: int,
        tenant_id: int = 1,
    ) -> Tuple[Optional[UserAgent], str]:
        existing_agents = self.db.query(UserAgent).filter(
            UserAgent.user_id == user_id,
            UserAgent.agent_status == 1
        ).all()
        
        if existing_agents:
            if len(existing_agents) == 1:
                return existing_agents[0], "获取已开通智能体"
            else:
                default_template = self._get_default_template()
                if default_template:
                    for agent in existing_agents:
                        if agent.template_id == default_template.id:
                            return agent, "获取默认模板智能体"
                return existing_agents[0], "获取首个已开通智能体"
        
        result, message = self.enable_agent_for_user(
            user_id=user_id,
            template_id=None,
            tenant_id=tenant_id,
        )
        
        if result:
            return result.user_agent, "自动开通智能体成功"
        else:
            return None, message
    
    def _get_conversation_requirements_data(self, conversation: AgentConversation) -> Dict[str, Any]:
        requirements = conversation.requirements_json
        if isinstance(requirements, dict):
            return requirements
        if isinstance(requirements, str) and requirements:
            try:
                return json.loads(requirements)
            except json.JSONDecodeError:
                return {}
        return {}
    
    def _get_welcome_message_from_template(self, user_agent: UserAgent) -> str:
        """
        从模板快照中获取欢迎语
        """
        default_welcome = """你好，我是企业 PPT 助手。

你可以直接告诉我：
1. 这份 PPT 是做什么用的
2. 是给谁看的
3. 想做多少页
4. 有没有必须引用的公司资料

如果你暂时说不全，我也会一步一步带你把这件事理顺。"""
        
        try:
            config_snapshot = user_agent.config_snapshot
            if config_snapshot and isinstance(config_snapshot, dict):
                welcome = config_snapshot.get("welcome_message")
                if welcome and isinstance(welcome, str) and len(welcome) > 0:
                    return welcome
        except Exception:
            pass
        
        return default_welcome
    
    def _is_ppt_assistant(self, user_agent: UserAgent) -> bool:
        """
        判断是否是 PPT 助手模板
        """
        try:
            config_snapshot = user_agent.config_snapshot
            if config_snapshot and isinstance(config_snapshot, dict):
                code = config_snapshot.get("code", "")
                if "ppt" in code.lower():
                    return True
        except Exception:
            pass
        
        return True
    
    def _get_knowledge_categories_from_template(self, user_agent: UserAgent) -> Optional[List[str]]:
        """
        从模板快照中获取知识分类配置
        """
        try:
            config_snapshot = user_agent.config_snapshot
            if config_snapshot and isinstance(config_snapshot, dict):
                knowledge_scope = config_snapshot.get("knowledge_scope", "none")
                knowledge_categories = config_snapshot.get("knowledge_categories")
                
                if knowledge_scope == "none":
                    return None
                elif knowledge_scope == "global":
                    return []
                elif knowledge_scope == "category" and knowledge_categories:
                    if isinstance(knowledge_categories, list):
                        return [str(c) for c in knowledge_categories if c]
        except Exception as e:
            print(f"获取知识分类配置失败: {e}")
        
        return None
    
    def _build_knowledge_query(self, requirements, message: str) -> str:
        """
        根据需求和消息构建知识检索查询
        """
        query_parts = []
        
        if requirements.topic:
            query_parts.append(requirements.topic)
        if requirements.audience:
            query_parts.append(requirements.audience)
        if requirements.scene:
            query_parts.append(requirements.scene)
        
        if message and len(message) > 3:
            query_parts.append(message)
        
        return " ".join(query_parts) if query_parts else ""
    
    def send_chat_message(
        self,
        agent_id: int,
        user_id: int,
        conversation_id: Optional[int],
        message: Optional[str],
        action_type: str = "message",
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Tuple[Optional[ChatResponse], str]:
        user_agent = self.get_user_agent_by_id(agent_id, user_id)
        if not user_agent:
            return None, "智能体不存在或无权限访问"
        
        if conversation_id:
            conversation = self.db.query(AgentConversation).filter(
                AgentConversation.id == conversation_id,
                AgentConversation.user_agent_id == agent_id,
                AgentConversation.user_id == user_id
            ).first()
            if not conversation:
                return None, "会话不存在或无权限访问"
        else:
            conversation = AgentConversation(
                user_id=user_id,
                user_agent_id=agent_id,
                hermes_profile=user_agent.hermes_profile,
                title=message[:50] if message else "新对话",
                current_stage="welcome",
                status=1,
                message_count=0,
                latest_user_input=message[:1000] if message else None,
                requirements_json={},
                structured_result_json=None,
                messages_json=[],
            )
            self.db.add(conversation)
            self.db.commit()
            self.db.refresh(conversation)
        
        user_agent.last_used_at = datetime.utcnow()
        conversation.last_message_at = datetime.utcnow()
        if message:
            conversation.latest_user_input = message[:1000]

        if self._is_ppt_assistant(user_agent) and message:
            conversation.requirements_json = self._extract_ppt_requirements_from_user_text(
                message,
                self._get_conversation_requirements_data(conversation),
            )
        
        if self._is_ppt_assistant(user_agent):
            try:
                if not message or not message.strip():
                    return None, "消息不能为空"

                action_message_map = {
                    "confirm_template": "我确认使用公司标准模板。",
                    "confirm_outline": "我确认每页内容草稿，可以生成。",
                    "confirm_generation": "确认生成 PPT 文件。",
                }
                effective_user_message = action_message_map.get(action_type, message)

                self._append_message(conversation, "user", effective_user_message)

                hermes_service = self._get_hermes_chat_service()
                planner_result = hermes_service.plan_ppt_flow(
                    user_message=effective_user_message,
                    existing_requirements=self._get_conversation_requirements_data(conversation),
                    current_stage=conversation.current_stage,
                    profile_name=conversation.hermes_profile or user_agent.hermes_profile,
                    hermes_session_id=conversation.hermes_conversation_id,
                )
                if planner_result.hermes_session_id and planner_result.hermes_session_id != conversation.hermes_conversation_id:
                    conversation.hermes_conversation_id = planner_result.hermes_session_id

                merged: Dict[str, Any] = dict(self._get_conversation_requirements_data(conversation) or {})
                if isinstance(planner_result.requirements, dict):
                    for k, v in planner_result.requirements.items():
                        if v is None:
                            continue
                        if isinstance(v, str) and not v.strip():
                            continue
                        merged[k] = v
                conversation.requirements_json = merged
                conversation.current_stage = planner_result.stage or conversation.current_stage

                assistant_content = ""
                structured_result = None

                if isinstance(planner_result.ppt_content, dict) and planner_result.ppt_content.get("type") == "ppt_content":
                    structured_result = planner_result.ppt_content
                    conversation.structured_result_json = structured_result
                    conversation.current_stage = "content_draft"
                    assistant_content = self._format_outline_markdown(structured_result) or "已生成每页内容草稿。"
                    assistant_content += "\n\n请确认是否可以生成（回复“确认”或点击「确认大纲」），或直接告诉我需要改哪一页。"
                    conversation.outline_confirmed = False
                    conversation.template_confirmed = bool(str(merged.get("template_path") or "").strip() or str(merged.get("style") or "").strip())
                    conversation.final_generation_confirmed = False

                else:
                    qs = planner_result.questions or []
                    if qs:
                        assistant_content = "\n".join([f"- {q}" for q in qs])
                    else:
                        assistant_content = "我已收到。你可以继续补充需求信息，或者确认模板/内容。"

                    if conversation.current_stage == "collecting_requirements":
                        conversation.outline_confirmed = False
                        conversation.template_confirmed = False
                        conversation.final_generation_confirmed = False
                    elif conversation.current_stage == "template_select":
                        conversation.outline_confirmed = False
                        conversation.template_confirmed = False
                        conversation.final_generation_confirmed = False
                    elif conversation.current_stage == "ready_generate":
                        conversation.outline_confirmed = True
                        conversation.template_confirmed = True
                        conversation.final_generation_confirmed = False
                        assistant_content = "已确认内容与模板，可以点击「生成 PPT」生成文件。"

                self._append_message(conversation, "assistant", assistant_content)
                self.db.commit()

                assistant_message = {"role": "assistant", "content": assistant_content}
                return ChatResponse(
                    conversation=conversation,
                    assistant_message=assistant_message,
                    structured_result=structured_result if isinstance(structured_result, dict) else (conversation.structured_result_json if isinstance(conversation.structured_result_json, dict) else None),
                ), "发送成功"

            except Exception as e:
                err_text = str(e)
                logger.error(f"Hermes service unavailable: {err_text}")
                if "No inference provider configured" in err_text or "AuthError" in err_text:
                    return None, (
                        "Hermes 未配置推理服务提供商或密钥。"
                        "请先在 Hermes 侧完成模型与密钥配置："
                        "运行 `hermes model` 选择 provider/model，"
                        "并在 `~/.hermes/.env`（或对应 Profile 的 .env）配置 API key，"
                        "然后重试。"
                    )
                return None, f"Hermes 调用失败：{err_text}"
        else:
            return None, "该智能体类型暂不支持"
    
    def get_conversations(
        self,
        agent_id: int,
        user_id: int,
        page: int = 1,
        page_size: int = 20,
        status: Optional[int] = None,
    ) -> AgentConversationListResponse:
        query = self.db.query(AgentConversation).filter(
            AgentConversation.user_agent_id == agent_id,
            AgentConversation.user_id == user_id
        )
        
        if status is not None:
            query = query.filter(AgentConversation.status == status)
        
        total = query.count()
        total_pages = math.ceil(total / page_size) if total > 0 else 1
        
        offset = (page - 1) * page_size
        conversations = query.order_by(
            AgentConversation.last_message_at.desc()
        ).offset(offset).limit(page_size).all()
        
        from app.schemas import AgentConversationResponse
        
        return AgentConversationListResponse(
            items=[AgentConversationResponse.model_validate(c) for c in conversations],
            total=total,
            page=page,
            page_size=page_size,
            total_pages=total_pages,
        )
    
    def get_conversation_detail(
        self,
        conversation_id: int,
        agent_id: int,
        user_id: int,
    ) -> Tuple[Optional[ConversationDetailResponse], str]:
        conversation = self.db.query(AgentConversation).filter(
            AgentConversation.id == conversation_id,
            AgentConversation.user_agent_id == agent_id,
            AgentConversation.user_id == user_id
        ).first()
        
        if not conversation:
            return None, "会话不存在或无权限访问"
        
        files = self.db.query(AgentGeneratedFile).filter(
            AgentGeneratedFile.conversation_id == conversation_id,
            AgentGeneratedFile.user_id == user_id,
            AgentGeneratedFile.deleted_at.is_(None)
        ).all()
        
        file_dicts = []
        for f in files:
            file_dicts.append({
                "id": f.id,
                "user_id": f.user_id,
                "user_agent_id": f.user_agent_id,
                "conversation_id": f.conversation_id,
                "file_type": f.file_type,
                "file_name": f.file_name,
                "file_path": f.file_path,
                "file_size": f.file_size,
                "mime_type": f.mime_type,
                "template_name": f.template_name,
                "source_type": f.source_type,
                "version_no": f.version_no,
                "generation_status": f.generation_status,
                "error_message": f.error_message,
                "created_at": f.created_at.isoformat() if f.created_at else None,
                "updated_at": f.updated_at.isoformat() if f.updated_at else None,
            })
        
        from app.schemas import AgentConversationResponse
        
        return ConversationDetailResponse(
            conversation=AgentConversationResponse.model_validate(conversation),
            messages=conversation.messages_json or [],
            structured_result=conversation.structured_result_json,
            files=file_dicts,
        ), "获取成功"
    
    def generate_ppt(
        self,
        agent_id: int,
        conversation_id: int,
        user_id: int,
        template_name: Optional[str],
        regenerate: bool = False,
    ) -> Tuple[Optional[AgentGeneratedFile], str]:
        user_agent = self.get_user_agent_by_id(agent_id, user_id)
        if not user_agent:
            return None, "智能体不存在或无权限访问"
        
        conversation = self.db.query(AgentConversation).filter(
            AgentConversation.id == conversation_id,
            AgentConversation.user_agent_id == agent_id,
            AgentConversation.user_id == user_id
        ).first()
        
        if not conversation:
            return None, "会话不存在或无权限访问"

        structured_result = conversation.structured_result_json
        slides_data = None
        if isinstance(structured_result, dict):
            slides = structured_result.get("slides")
            if isinstance(slides, list) and slides:
                slides_data = slides

        if not slides_data:
            try:
                requirements_data = self._get_conversation_requirements_data(conversation)
                missing = []
                missing.extend(self._get_missing_requirements_for_stage("collecting_requirements", requirements_data))
                missing.extend(self._get_missing_requirements_for_stage("template_select", requirements_data))
                missing = list(dict.fromkeys(missing))
                if missing:
                    return None, "生成 PPT 前仍缺少：" + "、".join(missing)

                generated = self._generate_ppt_structured_result(
                    user_agent=user_agent,
                    conversation=conversation,
                    requirements=requirements_data,
                )
                if generated:
                    conversation.structured_result_json = generated
                    conversation.current_stage = "content_ready"
                    self._append_message(conversation, "assistant", "已整理出可生成的 PPT 结构化内容，开始生成文件。")
                    self.db.commit()
                    structured_result = generated
                    slides = generated.get("slides") if isinstance(generated, dict) else None
                    if isinstance(slides, list) and slides:
                        slides_data = slides
            except Exception as e:
                logger.error(f"Failed to force structured ppt content from Hermes: {e}")

        if not slides_data:
            return None, "当前会话还没有可生成的 PPT 内容，请继续对话补齐信息后再生成"

        conversation.outline_confirmed = True
        conversation.template_confirmed = True
        conversation.final_generation_confirmed = True
        conversation.current_stage = "final_generating"
        
        existing_files = self.db.query(AgentGeneratedFile).filter(
            AgentGeneratedFile.conversation_id == conversation_id,
            AgentGeneratedFile.user_id == user_id,
            AgentGeneratedFile.deleted_at.is_(None)
        ).all()
        
        next_version = len(existing_files) + 1
        
        output_dir = os.path.join(
            settings.GENERATED_FILES_PATH,
            str(user_id),
            str(conversation_id)
        )
        os.makedirs(output_dir, exist_ok=True)
        
        file_name = f"PPT_{conversation_id}_v{next_version}.pptx"
        file_path = os.path.join(output_dir, file_name)
        
        knowledge_categories = self._get_knowledge_categories_from_template(user_agent)
        knowledge_context = None
        
        requirements_data = self._get_conversation_requirements_data(conversation)
        effective_template_name = template_name or requirements_data.get("style") or "默认模板"

        if knowledge_categories is not None:
            knowledge_retrieval_service = KnowledgeRetrievalService(self.db)
            topic = requirements_data.get("topic", "未设置主题")
            
            knowledge_context = knowledge_retrieval_service.get_knowledge_context(
                query=topic,
                categories=knowledge_categories if knowledge_categories else None,
                max_docs=3
            )
        
        pending_file = AgentGeneratedFile(
            user_id=user_id,
            user_agent_id=agent_id,
            conversation_id=conversation_id,
            file_type="pptx",
            file_name=file_name,
            file_path=file_path,
            file_size=0,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            template_name=effective_template_name,
            source_type="regenerated" if regenerate else "generated",
            version_no=next_version,
            generation_status=0,
        )
        
        self.db.add(pending_file)
        self.db.commit()
        self.db.refresh(pending_file)
        
        try:
            ppt_content = self._create_demo_ppt(
                conversation=conversation,
                version=next_version,
                template_name=effective_template_name,
                knowledge_context=knowledge_context
            )
            with open(file_path, "wb") as f:
                f.write(ppt_content)
            
            file_size = os.path.getsize(file_path)
            
            pending_file.file_size = file_size
            pending_file.generation_status = 1
            pending_file.error_message = None
            pending_file.updated_at = datetime.utcnow()
            
            self.db.commit()
            self.db.refresh(pending_file)
            
        except Exception as e:
            error_msg = str(e)
            pending_file.generation_status = 2
            pending_file.error_message = error_msg
            pending_file.updated_at = datetime.utcnow()
            self.db.commit()
            self.db.refresh(pending_file)
            
            self._log_operation(
                operator_user_id=user_id,
                target_type="file",
                target_id=pending_file.id,
                action="ppt:generate",
                action_name="生成PPT",
                success=False,
                details={
                    "conversation_id": conversation_id,
                    "template_name": effective_template_name,
                    "version_no": next_version,
                },
                error_message=error_msg,
            )
            
            return None, f"生成PPT失败: {error_msg}"
        
        conversation.final_file_id = pending_file.id
        conversation.current_stage = "completed"
        conversation.status = 2
        conversation.completed_at = datetime.utcnow()
        conversation.final_generation_confirmed = True
        conversation.messages_json = (conversation.messages_json or []) + [{
            "role": "assistant",
            "content": f"正式 PPT 已生成完成，文件名：{pending_file.file_name}",
        }]
        conversation.message_count = len(conversation.messages_json)
        
        self.db.commit()
        
        self._log_operation(
            operator_user_id=user_id,
            target_type="file",
            target_id=pending_file.id,
            action="ppt:generate",
            action_name="生成PPT",
            success=True,
            details={
                "conversation_id": conversation_id,
                "template_name": effective_template_name,
                "version_no": next_version,
            },
        )
        
        return pending_file, "PPT生成成功"
    
    def _create_demo_ppt(
        self, 
        conversation: AgentConversation, 
        version: int,
        template_name: Optional[str],
        knowledge_context: Optional[Dict[str, Any]] = None
    ) -> bytes:
        """
        创建一个演示用的 PPT 文件
        优先使用新的渲染服务，如果有结构化数据
        否则回退到旧的逻辑
        """
        try:
            requirements_data = self._get_conversation_requirements_data(conversation)
            structured_result = conversation.structured_result_json or {}
            
            topic = requirements_data.get("topic", "未设置主题")
            audience = requirements_data.get("audience")
            scene = requirements_data.get("scene")
            style = requirements_data.get("style")
            template_path = str(requirements_data.get("template_path") or "").strip()
            slides_data = structured_result.get("slides") if isinstance(structured_result, dict) else None
            
            if template_path and os.path.isfile(template_path):
                logger.info(f"Using template pptx as base: {template_path}")
                return self._create_demo_ppt_fallback(
                    conversation=conversation,
                    version=version,
                    template_name=template_name,
                    knowledge_context=knowledge_context,
                    template_path=template_path,
                )

            if slides_data and isinstance(slides_data, list) and len(slides_data) > 0:
                logger.info(f"Using new PPT renderer with {len(slides_data)} slides")
                
                effective_style = style or template_name or "默认"
                
                subtitle = ""
                if scene:
                    subtitle += f"场景：{scene}\n"
                if audience:
                    subtitle += f"受众：{audience}"
                
                metadata = {
                    "version": version,
                    "scene": scene,
                    "audience": audience,
                    "style": effective_style,
                }
                
                return create_ppt_from_data(
                    slides_data=slides_data,
                    title=topic,
                    subtitle=subtitle if subtitle else None,
                    style=effective_style,
                    metadata=metadata,
                )
            
            logger.info("No structured slide data found, using fallback renderer")
            return self._create_demo_ppt_fallback(
                conversation=conversation,
                version=version,
                template_name=template_name,
                knowledge_context=knowledge_context,
                template_path=None,
            )
            
        except ImportError:
            placeholder_content = self._create_placeholder_content(
                conversation=conversation,
                version=version,
                template_name=template_name,
                knowledge_context=knowledge_context,
            )
            return placeholder_content.encode('utf-8')
        except Exception as e:
            logger.error(f"PPT generation failed: {e}")
            placeholder_content = self._create_placeholder_content(
                conversation=conversation,
                version=version,
                template_name=template_name,
                knowledge_context=knowledge_context,
                error=str(e),
            )
            return placeholder_content.encode('utf-8')
    
    def _create_demo_ppt_fallback(
        self, 
        conversation: AgentConversation, 
        version: int,
        template_name: Optional[str],
        knowledge_context: Optional[Dict[str, Any]] = None,
        template_path: Optional[str] = None,
    ) -> bytes:
        """
        旧的 PPT 生成逻辑（作为 fallback）
        """
        from pptx import Presentation
        import io
        
        if template_path and os.path.isfile(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
        
        requirements_data = self._get_conversation_requirements_data(conversation)
        structured_result = conversation.structured_result_json or {}

        topic = requirements_data.get("topic", "未设置主题")
        audience = requirements_data.get("audience")
        scene = requirements_data.get("scene")
        page_count = requirements_data.get("page_count")
        style = requirements_data.get("style")
        slides_data = structured_result.get("slides") if isinstance(structured_result, dict) else None
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = topic
        subtitle_text = f"版本 v{version}\n模板: {template_name or '默认'}"
        if style:
            subtitle_text += f"\n风格: {style}"
        subtitle.text = subtitle_text
        
        outline_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(outline_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "目录"
        tf = body.text_frame
        tf.text = ""

        if audience or scene or page_count:
            p = tf.add_paragraph()
            p.text = "【基本信息】"
            p.bold = True

            if audience:
                p = tf.add_paragraph()
                p.text = f"• 目标受众：{audience}"
                p.level = 1
            if scene:
                p = tf.add_paragraph()
                p.text = f"• 使用场景：{scene}"
                p.level = 1
            if page_count:
                p = tf.add_paragraph()
                p.text = f"• 预计页数：{page_count}页"
                p.level = 1

        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = "【内容结构】"
        p.bold = True

        outline_titles = []
        if isinstance(slides_data, list) and slides_data:
            outline_titles = [f"• 第{slide_item.get('index', idx + 1)}页：{slide_item.get('title', '未命名页面')}" for idx, slide_item in enumerate(slides_data)]
        else:
            outline_titles = [
                "• 第1页：封面",
                "• 第2页：目录",
                "• 第3页：正文内容",
            ]

        for outline_title in outline_titles[: min(len(outline_titles), 10)]:
            p = tf.add_paragraph()
            p.text = outline_title
            p.level = 1

        if isinstance(slides_data, list) and slides_data:
            for slide_item in slides_data:
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]

                title.text = slide_item.get("title") or f"第{slide_item.get('index', '')}页"
                tf = body.text_frame
                subtitle = slide_item.get("subtitle")
                if subtitle:
                    tf.text = subtitle
                else:
                    tf.text = ""

                bullets = slide_item.get("bullets") or []
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0

                speaker_notes = slide_item.get("speaker_notes")
                if speaker_notes:
                    p = tf.add_paragraph()
                    p.text = ""
                    p = tf.add_paragraph()
                    p.text = f"演讲备注：{speaker_notes}"
                    p.level = 0
        else:
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = "正文内容"
            tf = body.text_frame
            tf.text = "这是一个由 AgentNow 智能体平台生成的 PPT。"

            if audience:
                p = tf.add_paragraph()
                p.text = f"• 目标受众：{audience}"
            if scene:
                p = tf.add_paragraph()
                p.text = f"• 使用场景：{scene}"
            if page_count:
                p = tf.add_paragraph()
                p.text = f"• 预计页数：{page_count}"
            if style:
                p = tf.add_paragraph()
                p.text = f"• 展示风格：{style}"
        
        if knowledge_context and knowledge_context.get("has_knowledge"):
            sources = knowledge_context.get("sources", [])
            if sources:
                reference_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(reference_slide_layout)
                title = slide.shapes.title
                body = slide.placeholders[1]
                
                title.text = "参考资料（企业知识库）"
                tf = body.text_frame
                tf.text = "以下内容来自企业知识库："
                
                for i, source in enumerate(sources, 1):
                    p = tf.add_paragraph()
                    p.text = ""
                    p = tf.add_paragraph()
                    p.text = f"{i}. 《{source['title']}》"
                    p.bold = True
                    if source.get("category"):
                        p = tf.add_paragraph()
                        p.text = f"   分类：{source['category']}"
                
                knowledge_text = knowledge_context.get("knowledge_text", "")
                if knowledge_text and len(knowledge_text) > 0:
                    p = tf.add_paragraph()
                    p.text = ""
                    p = tf.add_paragraph()
                    p.text = "【资料摘要】"
                    p.bold = True
                    
                    content_lines = knowledge_text[:800].split('\n')
                    for line in content_lines[:15]:
                        if line.strip():
                            p = tf.add_paragraph()
                            p.text = line[:80]
                            p.level = 1
        
        summary_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(summary_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = "总结与说明"
        tf = body.text_frame
        tf.text = "PPT 生成完成！"
        
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = "【生成信息】"
        p.bold = True
        
        p = tf.add_paragraph()
        p.text = f"• 会话ID: {conversation.id}"
        p = tf.add_paragraph()
        p.text = f"• 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        p = tf.add_paragraph()
        p.text = f"• 版本: v{version}"
        
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = "【说明】"
        p.bold = True
        
        if knowledge_context and knowledge_context.get("has_knowledge"):
            p = tf.add_paragraph()
            p.text = "✓ 本 PPT 内容参考了企业知识库中的相关资料"
            p.level = 1
        else:
            p = tf.add_paragraph()
            p.text = "⚠ 本 PPT 为通用结构建议，未找到企业知识库中的相关资料"
            p.level = 1
            p = tf.add_paragraph()
            p.text = "  如需更贴合企业实际的内容，请上传相关资料到知识库"
            p.level = 1
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.read()
    
    def _create_placeholder_content(
        self,
        conversation: AgentConversation,
        version: int,
        template_name: Optional[str],
        knowledge_context: Optional[Dict[str, Any]] = None,
        error: Optional[str] = None,
    ) -> str:
        """
        创建占位符内容（当 python-pptx 不可用时）
        """
        requirements_data = self._get_conversation_requirements_data(conversation)
        
        placeholder_content = f"""PPT 演示文件
================
会话ID: {conversation.id}
版本: v{version}
模板: {template_name or '默认'}
生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

说明：
- 此为演示版本
- 正式版本将生成真实的 .pptx 文件
- 请安装 python-pptx 库以启用完整功能
"""
        if error:
            placeholder_content += f"\n错误信息: {error}\n"
        
        placeholder_content += f"""
需求信息：
{json.dumps(requirements_data, ensure_ascii=False, indent=2) if requirements_data else '未记录详细需求'}

知识检索状态：
"""
        
        if knowledge_context and knowledge_context.get("has_knowledge"):
            placeholder_content += "✓ 已检索到企业知识库资料\n"
            sources = knowledge_context.get("sources", [])
            for i, source in enumerate(sources, 1):
                placeholder_content += f"  {i}. {source['title']}"
                if source.get("category"):
                    placeholder_content += f"（{source['category']}）"
                placeholder_content += "\n"
        else:
            placeholder_content += "⚠ 未检索到企业知识库中的相关资料\n"
            placeholder_content += "  本内容为通用建议\n"
        
        return placeholder_content
    
    def download_generated_file(
        self,
        file_id: int,
        agent_id: int,
        user_id: int,
    ) -> Tuple[Optional[bytes], Optional[str], str, str]:
        generated_file = self.db.query(AgentGeneratedFile).filter(
            AgentGeneratedFile.id == file_id,
            AgentGeneratedFile.user_agent_id == agent_id,
            AgentGeneratedFile.user_id == user_id,
            AgentGeneratedFile.deleted_at.is_(None)
        ).first()
        
        if not generated_file:
            return None, None, "", "文件不存在或无权限访问"
        
        if generated_file.generation_status != 1:
            return None, None, "", "文件生成未完成或已失败"
        
        file_path = generated_file.file_path
        if not file_path or not os.path.exists(file_path):
            return None, None, "", "文件已删除或不可访问"
        
        try:
            with open(file_path, "rb") as f:
                file_content = f.read()
            return file_content, generated_file.mime_type, generated_file.file_name, "获取成功"
        except Exception as e:
            return None, None, "", f"读取文件失败: {str(e)}"
